import os
from openai import OpenAI
from langchain.text_splitter import RecursiveCharacterTextSplitter
from sentence_transformers import SentenceTransformer
from supabase import create_client
from flask import current_app
from utils.supabase_client import get_supabase

from docx import Document
import fitz
from pptx import Presentation
import win32com.client


def extract_text_from_file(file_path: str) -> str:
    """
    根据文件扩展名提取文本内容
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == '.txt':
            return extract_txt(file_path)
        elif ext == '.docx':
            return extract_docx(file_path)
        elif ext == '.doc':
            return extract_doc(file_path)
        elif ext == '.pdf':
            return extract_pdf(file_path)
        elif ext == '.pptx':
            return extract_powerpoint(file_path)
        else:
            raise ValueError(f"不支持的文件格式: {ext}")
    except Exception as e:
        raise Exception(f"提取文本失败: {str(e)}")

def extract_txt(file_path: str) -> str:
    """提取txt文件文本"""
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

def extract_docx(file_path: str) -> str:
    """提取docx文件文本"""
    doc = Document(file_path)
    text = []
    for paragraph in doc.paragraphs:
        text.append(paragraph.text)
    return '\n'.join(text)

def extract_doc(file_path: str) -> str:
    """提取doc文件文本"""
    try:
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(os.path.abspath(file_path))
        text = doc.Content.Text
        doc.Close()
        word.Quit()
        return text
    except Exception as e:
        print(f"提取失败: {e}")

def extract_pdf(file_path: str) -> str:
    """提取PDF文件文本"""
    text = []
    
    # 使用PyMuPDF
    try:
        doc = fitz.open(file_path)
        for page in doc:
            page_text = page.get_text()
            if page_text.strip():
                text.append(page_text)
        doc.close()
        if text:
            return '\n'.join(text)
    except Exception as e:
        raise Exception(f"PDF文本提取失败: {str(e)}")


def extract_powerpoint(file_path: str) -> str:
    """提取PowerPoint文件文本"""
    prs = Presentation(file_path)
    text = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        # text.append(f"幻灯片 {slide_num}:")
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text)
        text.append('')  # 空行分隔
    
    return '\n'.join(text)



# === 加载中文向量模型 ===
model = SentenceTransformer("BAAI/bge-small-zh")

# ==== 提取向量 ====
def get_embedding(text):
    return model.encode(text, normalize_embeddings=True).tolist()


# ==== 向量化文件 ====
def embed_file(file_path: str, file_id: int):
    supabase = get_supabase()
    # 👇 获取 user_id
    file_record = supabase.table("files").select("user_id").eq("id", file_id).single().execute()
    user_id = file_record.data["user_id"]
    # with open(file_path, 'r', encoding='utf-8') as f:
    #     text = f.read()
    # 提取文本内容
    try:
        text = extract_text_from_file(file_path)
        if not text.strip():
            raise Exception("提取的文本内容为空")
    except Exception as e:
        raise Exception(f"文本提取失败: {str(e)}")

    splitter = RecursiveCharacterTextSplitter(
        separators=["\n\n", "\n", "。", "！", "？", "，", " ", ""],
        chunk_size=500, #每块文本最多500个字符
        chunk_overlap=50 #相邻块之间有50个字符的重叠（避免重要信息在分割点丢失）
    )
    chunks = splitter.split_text(text)


    for i, chunk in enumerate(chunks):
        try:
            embedding = get_embedding(chunk)
            supabase.table("user_vectors").insert({
                "file_id": file_id,  # ✅ 用外键指向 files 表
                "user_id": user_id,
                "content": chunk,
                "embedding": embedding,
                "metadata": {
                    "chunk_index": i,
                    "file_type": os.path.splitext(file_path)[1].lower()
                }
            }).execute()
            print(f"[✅] 插入 file_id={file_id} 第 {i + 1} 段")
        except Exception as e:
            print(f"[❌] 插入失败（file_id={file_id}, chunk={i + 1}）：{e}")
